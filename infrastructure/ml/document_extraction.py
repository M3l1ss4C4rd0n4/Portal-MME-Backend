"""
Extracción de texto de documentos PDF/PPTX para RAG (ontologia.informes_texto_embeddings).

Chunking natural: 1 chunk = 1 página (PDF) o 1 diapositiva (PPTX) — evita la
complejidad de un splitter genérico de tokens para un corpus inicial pequeño
(~10 documentos). Cada chunk conserva su número de página/diapositiva (1-indexed)
como chunk_index, útil para citar la fuente exacta en una respuesta del chatbot.
"""

from typing import List, Optional

from infrastructure.logging.logger import get_logger

logger = get_logger(__name__)

# Subido de 30 a 80 (Fase 23, 2026-08-06): con 30, páginas-título casi vacías
# (ej. "Informe diario de situación energética", 38 caracteres) pasaban el
# filtro y embebían con similitud engañosamente alta frente a consultas de
# dominio reales — reproducido en vivo devolviendo la misma frase 5 veces en
# el top-5 de una búsqueda sobre un tema completamente distinto.
LONGITUD_MINIMA_CHUNK = 80


def extraer_paginas_pdf(contenido: bytes, max_paginas: Optional[int] = None) -> List[str]:
    """Texto por página de un PDF. Páginas con menos de LONGITUD_MINIMA_CHUNK
    caracteres (portadas, separadores) se omiten — no aportan a la búsqueda.

    `max_paginas`: opcional, limita cuántas páginas iniciales se procesan.
    Ninguna fuente existente lo pasa (comportamiento por defecto sin cambios).
    Agregado para INFODESPACHO de XM (build_informes_embeddings.py) — informes
    diarios de 27-35 MB / 25-30 páginas donde, verificado en 2 muestras reales,
    solo las primeras 3 páginas son narrativas (Resumen de generación,
    Novedades, Consideración/Indisponibilidades); el resto son tablas horarias
    de generación por planta que no aportan al RAG y solo agregarían ruido."""
    import pdfplumber
    import io

    paginas: List[str] = []
    with pdfplumber.open(io.BytesIO(contenido)) as pdf:
        pages = pdf.pages[:max_paginas] if max_paginas else pdf.pages
        for page in pages:
            texto = (page.extract_text() or "").strip()
            if len(texto) >= LONGITUD_MINIMA_CHUNK:
                paginas.append(texto)
    return paginas


def extraer_slides_pptx(contenido: bytes) -> List[str]:
    """
    Texto por diapositiva de un PPTX: concatena todos los text frames (títulos,
    cuerpos, tablas) más las notas del orador. Gráficos/imágenes sin texto
    vectorial no se extraen (no hay OCR) — limitación conocida, no todo el
    contenido visual de una diapositiva queda indexado.
    """
    from pptx import Presentation
    import io

    slides: List[str] = []
    prs = Presentation(io.BytesIO(contenido))
    for slide in prs.slides:
        partes: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                partes.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    fila = " | ".join(cell.text.strip() for cell in row.cells)
                    if fila.strip(" |"):
                        partes.append(fila)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notas = slide.notes_slide.notes_text_frame.text.strip()
            if notas:
                partes.append(notas)
        texto = "\n".join(partes).strip()
        if len(texto) >= LONGITUD_MINIMA_CHUNK:
            slides.append(texto)
    return slides


def extraer_parrafos_docx(contenido: bytes) -> List[str]:
    """
    Texto de un DOCX agrupado en chunks de ~15 párrafos (un DOCX no tiene páginas
    reales hasta que se renderiza — a diferencia de PDF/PPTX no hay una unidad
    natural de página/diapositiva, así que se agrupa por cantidad de párrafos).
    """
    from docx import Document
    import io

    PARRAFOS_POR_CHUNK = 15
    doc = Document(io.BytesIO(contenido))
    parrafos = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    chunks: List[str] = []
    for i in range(0, len(parrafos), PARRAFOS_POR_CHUNK):
        texto = "\n".join(parrafos[i : i + PARRAFOS_POR_CHUNK])
        if len(texto) >= LONGITUD_MINIMA_CHUNK:
            chunks.append(texto)
    return chunks


def chunk_texto_plano(texto: str, tam_objetivo: int = 2000) -> List[str]:
    """
    Agrupa texto plano ya extraído (no un archivo binario — ej. HTML→texto de
    una resolución CREG, sin páginas/diapositivas naturales) en chunks por
    párrafo, acumulando hasta acercarse a `tam_objetivo` caracteres antes de
    cortar — evita partir un párrafo/artículo a la mitad como haría un corte
    fijo por caracteres. Usado por build_informes_embeddings.py::_indexar_creg_normativa.
    """
    parrafos = [p.strip() for p in texto.split("\n") if p.strip()]
    chunks: List[str] = []
    actual: List[str] = []
    largo_actual = 0
    for parrafo in parrafos:
        if largo_actual + len(parrafo) > tam_objetivo and actual:
            chunks.append("\n".join(actual))
            actual, largo_actual = [], 0
        actual.append(parrafo)
        largo_actual += len(parrafo)
    if actual:
        chunks.append("\n".join(actual))
    return [c for c in chunks if len(c) >= LONGITUD_MINIMA_CHUNK]


def extraer_chunks(
    contenido: bytes, tipo_archivo: str, max_paginas: Optional[int] = None
) -> List[str]:
    """Despacha a la función de extracción según tipo_archivo ('pdf' | 'pptx' | 'docx').
    `max_paginas` solo aplica a 'pdf' (ver extraer_paginas_pdf)."""
    if tipo_archivo == "pdf":
        return extraer_paginas_pdf(contenido, max_paginas=max_paginas)
    if tipo_archivo == "pptx":
        return extraer_slides_pptx(contenido)
    if tipo_archivo == "docx":
        return extraer_parrafos_docx(contenido)
    raise ValueError(f"Tipo de archivo no soportado para extracción: {tipo_archivo!r}")
